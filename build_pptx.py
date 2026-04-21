"""Generate rxcore.pptx — RxCore & KCS customer-facing deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
BG       = RGBColor(0x16, 0x1b, 0x27)
BG_CARD  = RGBColor(0x1c, 0x22, 0x35)
BG_DARK  = RGBColor(0x0f, 0x11, 0x17)
BLUE     = RGBColor(0x4f, 0x9c, 0xf9)
PURPLE   = RGBColor(0xa7, 0x8b, 0xfa)
GREEN    = RGBColor(0x34, 0xd3, 0x99)
AMBER    = RGBColor(0xfb, 0xbf, 0x24)
RED      = RGBColor(0xf8, 0x71, 0x71)
INDIGO   = RGBColor(0x81, 0x8c, 0xf8)
LIME     = RGBColor(0xa3, 0xe6, 0x35)
ORANGE   = RGBColor(0xfb, 0x92, 0x3c)
ROSE     = RGBColor(0xfb, 0x71, 0x85)
WHITE    = RGBColor(0xff, 0xff, 0xff)
GREY_LT  = RGBColor(0xc5, 0xcd, 0xe0)
GREY_MD  = RGBColor(0x9a, 0xa5, 0xbc)
GREY_DK  = RGBColor(0x4a, 0x57, 0x6e)
BORDER   = RGBColor(0x25, 0x2f, 0x44)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0.75)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, size=Pt(11), color=WHITE, bold=False,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def footer(slide, label, counter):
    add_rect(slide, Inches(0), H - Inches(0.45), W, Inches(0.45),
             fill=RGBColor(0x13, 0x17, 0x22), line_color=BORDER)
    add_textbox(slide, Inches(0.4), H - Inches(0.42), Inches(6), Inches(0.42),
                label.upper(), size=Pt(9), color=GREY_DK, bold=True)
    add_textbox(slide, W - Inches(2.5), H - Inches(0.42), Inches(1.2), Inches(0.42),
                counter, size=Pt(9), color=RGBColor(0x6b, 0x7a, 0x93), align=PP_ALIGN.RIGHT)
    add_textbox(slide, W - Inches(1.2), H - Inches(0.42), Inches(0.9), Inches(0.42),
                "Reshape", size=Pt(11), color=BLUE, bold=True, align=PP_ALIGN.RIGHT)

def header_block(slide, title_plain, title_accent, subtitle, tag):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.6), fill=RGBColor(0x18, 0x1f, 0x32))
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.16), Inches(9.5), Inches(0.72))
    tf = txBox.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    if title_plain:
        r1 = p.add_run(); r1.text = title_plain + " "
        r1.font.size = Pt(24); r1.font.bold = True; r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = title_accent
    r2.font.size = Pt(24); r2.font.bold = True; r2.font.color.rgb = BLUE
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.88), Inches(10.8), Inches(0.62))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r = p2.add_run(); r.text = subtitle
    r.font.size = Pt(11); r.font.color.rgb = GREY_MD
    add_textbox(slide, Inches(11.0), Inches(0.22), Inches(2.1), Inches(0.4),
                tag.upper(), size=Pt(8), color=GREY_DK, bold=True, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title, label, bullets, icon="", accent=BLUE):
    add_rect(slide, x, y, w, h, fill=BG_CARD, line_color=BORDER)
    icon_sz = Inches(0.34)
    add_rect(slide, x + Inches(0.16), y + Inches(0.16), icon_sz, icon_sz,
             fill=RGBColor(0x1a, 0x3a, 0x6b), line_color=None)
    add_textbox(slide, x + Inches(0.14), y + Inches(0.13), icon_sz + Inches(0.04), icon_sz,
                icon, size=Pt(14), align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.62), y + Inches(0.16), w - Inches(0.78), Inches(0.3),
                title, size=Pt(11), color=GREY_LT, bold=True)
    add_textbox(slide, x + Inches(0.62), y + Inches(0.44), w - Inches(0.78), Inches(0.22),
                label.upper(), size=Pt(8), color=accent, bold=True)
    ty = y + Inches(0.74)
    for b in bullets:
        dot = slide.shapes.add_shape(9, x + Inches(0.20), ty + Inches(0.10),
                                     Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        tb = slide.shapes.add_textbox(x + Inches(0.36), ty, w - Inches(0.52), Inches(0.56))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        if ' — ' in b:
            head, rest = b.split(' — ', 1)
            r1 = p.add_run(); r1.text = head
            r1.font.bold = True; r1.font.size = Pt(10); r1.font.color.rgb = GREY_LT
            r2 = p.add_run(); r2.text = ' — ' + rest
            r2.font.size = Pt(10); r2.font.color.rgb = GREY_MD
        else:
            r = p.add_run(); r.text = b
            r.font.size = Pt(10); r.font.color.rgb = GREY_MD
        ty += Inches(0.60)

def pipe_box(slide, x, y, w, h, label, title, lines):
    add_rect(slide, x, y, w, h, fill=BG_CARD, line_color=BORDER)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.2), Inches(0.2),
                label.upper(), size=Pt(8), color=BLUE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.30), w - Inches(0.2), Inches(0.25),
                title, size=Pt(11), color=GREY_LT, bold=True)
    ty = y + Inches(0.62)
    for ln in lines:
        tb = slide.shapes.add_textbox(x + Inches(0.15), ty, w - Inches(0.25), Inches(0.36))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        if ln.startswith("•") and ' — ' in ln:
            head, rest = ln[2:].split(' — ', 1)
            r1 = p.add_run(); r1.text = "• " + head
            r1.font.bold = True; r1.font.size = Pt(10); r1.font.color.rgb = GREY_LT
            r2 = p.add_run(); r2.text = ' — ' + rest
            r2.font.size = Pt(10); r2.font.color.rgb = GREY_MD
        else:
            r = p.add_run(); r.text = ln
            r.font.size = Pt(10); r.font.color.rgb = GREY_MD
        ty += Inches(0.34)

# ── SLIDE 1: COVER ─────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1, RGBColor(0x14, 0x1a, 0x28))
add_rect(s1, Inches(2), Inches(0), Inches(9.3), Inches(2.6),
         fill=RGBColor(0x1a, 0x2a, 0x45), line_color=None)
add_textbox(s1, Inches(1), Inches(1.1), Inches(11.3), Inches(0.4),
            "RESHAPE PLATFORM", size=Pt(11), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
txBox = s1.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.3), Inches(1.8))
tf = txBox.text_frame; tf.word_wrap = False
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "Industrial AI that knows "
r1.font.size = Pt(40); r1.font.bold = True; r1.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "your catalog"
r2.font.size = Pt(40); r2.font.bold = True; r2.font.color.rgb = BLUE
add_textbox(s1, Inches(1.5), Inches(3.55), Inches(10.3), Inches(1.0),
            "RxCore is the intelligence engine behind Reshape's agents. Trained on your product data, connected to your systems, and built to answer industrial automation questions with near-perfect accuracy — at scale.",
            size=Pt(13), color=GREY_MD, align=PP_ALIGN.CENTER)
chips = ["99.9% accuracy target", "API · MCP · A2A ready", "Siemens · Turck · Festo & more"]
cx = Inches(1.5)
for chip in chips:
    cw = Inches(3.3)
    add_rect(s1, cx, Inches(4.75), cw, Inches(0.50),
             fill=RGBColor(0x1e, 0x26, 0x36), line_color=BORDER)
    add_textbox(s1, cx, Inches(4.75), cw, Inches(0.50),
                chip, size=Pt(11), color=GREY_LT, bold=True, align=PP_ALIGN.CENTER)
    cx += cw + Inches(0.47)
footer(s1, "RxCore & KCS — Platform Overview", "1 / 6")

# ── SLIDE 2: WHAT IS RXCORE ────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2)
header_block(s2, "RxCore is a", "platform, not a chatbot",
             "Think of it as the engine under the hood. Your customer-facing tools and team applications sit on top of it — RxCore provides the intelligence, data, and accuracy layer that makes all of them reliable.",
             "Slide 2 · Platform Overview")
layers = [
    ("YOUR APPLICATIONS",        WHITE,  "Chatbots, quote builders, configurators, search interfaces — your UX and business logic, built on top of the core"),
    ("YOUR EXISTING SYSTEMS",    WHITE,  "ERP, CRM, website, email — RxCore connects to what you already have, without replacing it"),
    ("GROUNDING LAYER",          PURPLE, "The reasoning engine — understands PLCs, SCADA, drives, sensors, and 100,000+ SKUs. Everything between the user's question and an accurate answer."),
    ("RXCORE PLATFORM",          BLUE,   "Centralized intelligence & infrastructure — accessible via API, tool, MCP, or agent-to-agent (A2A). Everything above depends on it."),
    ("KNOWLEDGE CONSTRUCTION",   WHITE,  "Pre-builds your product knowledge — catalog, specs, compatibility, relationships — so agents answer instantly and reliably"),
]
ly = Inches(1.72)
lh = Inches(0.76)
for name, accent, desc in layers:
    bg = BG_CARD; lc = BORDER
    if accent == BLUE:   bg = RGBColor(0x1a, 0x2a, 0x45); lc = BLUE
    elif accent == PURPLE: bg = RGBColor(0x22, 0x1a, 0x3d); lc = PURPLE
    add_rect(s2, Inches(0.5), ly, Inches(12.33), lh, fill=bg, line_color=lc)
    add_textbox(s2, Inches(0.7), ly + Inches(0.08), Inches(2.4), Inches(0.28),
                name, size=Pt(9), color=accent, bold=True)
    add_textbox(s2, Inches(3.2), ly + Inches(0.08), Inches(9.4), Inches(0.62),
                desc, size=Pt(11), color=GREY_MD)
    ly += lh + Inches(0.07)
footer(s2, "Platform Architecture", "2 / 6")

# ── SLIDE 3: THREE PILLARS ─────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3)
header_block(s3, "How the grounding layer", "works",
             "Three engineering disciplines work together to turn a general-purpose AI into a specialist in your product catalog and industrial domain.",
             "Slide 3 · How It Works")
cw3 = Inches(4.0); ch3 = Inches(5.1); cy3 = Inches(1.72); cx3 = Inches(0.5)
pillars = [
    ("How it reasons", "Orchestration", "🎛️", BLUE, [
        "Purpose-built workflows — find a substitute, check compatibility, configure a product follow defined steps, not improvisation",
        "Right workflow, every time — the agent selects the appropriate pattern at the start of each turn",
        "Specialized sub-agents — complex queries split across focused agents that run in parallel",
        "Built-in safety checks — inputs and outputs validated automatically; unsafe requests blocked",
    ]),
    ("What it knows", "Knowledge Tools (from KCS)", "🔌", GREEN, [
        "Built from your catalog — queries structured tools generated from your own product data, not live brand APIs",
        "Understands relationships — what replaces what, what works with what, which families are compatible",
        "Lifecycle-aware — always checks active vs. end-of-life status before recommending any product",
        "Live inventory & pricing — stock and prices synced continuously so agents always quote accurately",
    ]),
    ("What it remembers", "Context Engineering", "🧠", RGBColor(0x2d, 0xd4, 0xbf), [
        "Three memory layers — long-term knowledge, the current conversation, and a fast-access cache",
        "No stale distraction — old tool results pruned so they don't mislead the model on later turns",
        "Smart summarization — long conversations compressed without losing key earlier context",
        "Curated, not accumulated — every token in context is there because it's useful",
    ]),
]
for title, label, icon, accent, bullets in pillars:
    card(s3, cx3, cy3, cw3, ch3, title, label, bullets, icon=icon, accent=accent)
    cx3 += cw3 + Inches(0.17)
footer(s3, "Reasoning · Knowledge · Memory", "3 / 6")

# ── SLIDE 4: PARALLEL RESEARCH + CONFIDENCE ────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4)
header_block(s4, "Before answering,", "RxCore checks everything at once",
             "The agent fans out across your entire knowledge layer simultaneously, scores every piece of evidence, and only composes an answer once all research is complete — never before.",
             "Slide 4 · Parallel Research & Confidence")
pw = Inches(3.6); ph = Inches(2.55); py4 = Inches(1.72)
arrow_w = Inches(0.44)
total = pw * 3 + arrow_w * 2
px4 = (W - total) / 2
pipe_box(s4, px4, py4, pw, ph, "Step 1", "Parallel fan-out", [
    "• Knowledge graph — relationships & inferred compatibility",
    "• Product database — official specs, lifecycle, replacements",
    "• Semantic search — fuzzy natural-language matching",
    "• Web fallback — fills gaps when internal data is incomplete",
])
add_textbox(s4, px4 + pw + Inches(0.05), py4 + Inches(1.0), arrow_w, Inches(0.5),
            "→", size=Pt(24), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
pipe_box(s4, px4 + pw + arrow_w, py4, pw, ph, "Step 2", "Evidence scoring", [
    "• Graph: edge weight & path length",
    "• Database: exact match vs. inferred join",
    "• Semantic: similarity score & relevance",
    '• Missing data = explicit "not found" — never silent',
])
add_textbox(s4, px4 + pw * 2 + arrow_w + Inches(0.05), py4 + Inches(1.0), arrow_w, Inches(0.5),
            "→", size=Pt(24), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
pipe_box(s4, px4 + pw * 2 + arrow_w * 2, py4, pw, ph, "Step 3", "Reconcile & tier", [
    "HIGH — Graph + DB agree → stated plainly",
    "MEDIUM — Single authoritative source → hedged",
    "LOW — Semantic only → uncertainty disclosed",
    "CONFLICT — Sources disagree → both surfaced",
])
cw4 = Inches(6.1); ch4 = Inches(2.55); cy4 = Inches(4.45)
card(s4, Inches(0.5), cy4, cw4, ch4, "Why this produces better answers",
     "No early commitment, no guessing", [
         "All research completes before any answer is written — agent cannot short-circuit with partial result",
         "Multi-source triangulation eliminates the most common AI failure: anchoring on first result",
         '"Not found" signals trigger escalation to broader search — never a fabricated answer',
     ], icon="⚡", accent=INDIGO)
card(s4, Inches(0.5) + cw4 + Inches(0.16), cy4, cw4, ch4, "What you see as a result",
     "Transparent, trustworthy answers", [
         "Every claim is sourced — official specs vs. inferred compatibility, clearly labelled",
         "Uncertainty is surfaced, not hidden — weak evidence disclosed, not masked with confidence",
         "Conflicts shown — when sources disagree, both perspectives surfaced for your team to decide",
     ], icon="🎯", accent=LIME)
footer(s4, "Parallel research · Scored evidence · Transparent answers", "4 / 6")

# ── SLIDE 5: KCS ──────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5)
header_block(s5, "We train a specialist", "before they start the job",
             "The Knowledge Construction System (KCS) pre-builds a structured, queryable knowledge layer from your catalog data — so the agent never scrambles at the moment a customer asks a question.",
             "Slide 5 · Knowledge Construction System")
pw5 = Inches(3.6); ph5 = Inches(2.3); py5 = Inches(1.72)
px5 = (W - (pw5 * 3 + arrow_w * 2)) / 2
pipe_box(s5, px5, py5, pw5, ph5, "Input", "Your data sources", [
    "• Files — PDF catalogs, datasheets, Excel specs",
    "• Web — brand websites, product pages, documentation",
    "• APIs — official brand APIs (e.g. Siemens SiePortal)",
    "• Experts — your team contributes knowledge via chat",
])
add_textbox(s5, px5 + pw5 + Inches(0.05), py5 + Inches(0.9), arrow_w, Inches(0.5),
            "→", size=Pt(24), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
pipe_box(s5, px5 + pw5 + arrow_w, py5, pw5, ph5, "Process", "KG Constructor agent", [
    "• Ingestion — parse, normalize, deduplicate all sources",
    "• Schema modeling — learns how this brand organizes products",
    "• Validation loop — flags conflicts, proposes rules",
    "• FDE review — your team approves structural changes",
])
add_textbox(s5, px5 + pw5 * 2 + arrow_w + Inches(0.05), py5 + Inches(0.9), arrow_w, Inches(0.5),
            "→", size=Pt(24), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
pipe_box(s5, px5 + pw5 * 2 + arrow_w * 2, py5, pw5, ph5, "Output", "Structured knowledge layer", [
    "• Knowledge graph — relationships, compatibility, replacements",
    "• Product database — specs, lifecycle, stock, pricing",
    "• Semantic index — natural-language product search",
])
cw5 = Inches(6.1); ch5 = Inches(2.7); cy5 = Inches(4.15)
card(s5, Inches(0.5), cy5, cw5, ch5, "Gets smarter with every brand",
     "Brand-specific learning", [
         "Memory across onboardings — lessons from Siemens inform how ABB or Rockwell is processed next",
         "Faster with each brand — 20 iterations for brand 1, 5 for brand 10",
         "Adapts to each catalog structure — Siemens, Rockwell, Schneider all organize products differently",
     ], icon="📈", accent=ORANGE)
card(s5, Inches(0.5) + cw5 + Inches(0.16), cy5, cw5, ch5, "Stays up to date automatically",
     "Continuous sync", [
         "Triggered by changes — new catalog files, API updates, or scheduled syncs kick off a refresh",
         "Low-risk changes auto-applied — price and stock go straight in; structural changes queued",
         "High-frequency pricing & inventory — synced more often than taxonomy",
     ], icon="🔄", accent=ROSE)
footer(s5, "Build-time knowledge · Brand learning · Continuous sync", "5 / 6")

# ── SLIDE 6: KNOWLEDGE LAYER → TOOLS ─────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6)
header_block(s6, "Three ways to answer", "any product question accurately",
             "KCS produces three complementary knowledge surfaces. The grounding layer queries all three through stable, generated tools — it never calls brand APIs at the moment a question is asked.",
             "Slide 6 · Knowledge Layer & Tools")
cw6 = Inches(4.0); ch6 = Inches(3.2); cy6 = Inches(1.72); cx6 = Inches(0.5)
surfaces = [
    ("Knowledge Graph", "Relationships & Inferred Compatibility", "🕸️", PURPLE,
     "Captures your product catalog structure and the relationships your data implies — even when no source explicitly states them. Powers 'what replaces this?' and cross-family recommendations.",
     ["graph_search", "get_taxonomy", "find_spare_parts", "recommend_modules"]),
    ("Product Database", "Official Source of Truth", "🗄️", GREEN,
     "What brand APIs and official sources explicitly declare: specs, certified replacements, lifecycle status, pricing, and stock. Deliberately separate so the agent is always transparent about its source.",
     ["sql_query", "get_stock_price", "get_alternatives", "check_compatibility"]),
    ("Semantic Index", "Natural Language Product Search", "✨", AMBER,
     "Embeddings from both the database and knowledge graph. Handles vague queries like 'something like the S7-1200 but cheaper' — surfaces candidates the agent enriches with structured data.",
     ["rag_search"]),
]
for title, label, icon, accent, desc, tools in surfaces:
    add_rect(s6, cx6, cy6, cw6, ch6, fill=BG_CARD, line_color=BORDER)
    add_rect(s6, cx6 + Inches(0.16), cy6 + Inches(0.16), Inches(0.34), Inches(0.34),
             fill=RGBColor(0x1a, 0x1a, 0x4e), line_color=None)
    add_textbox(s6, cx6 + Inches(0.14), cy6 + Inches(0.13), Inches(0.38), Inches(0.38),
                icon, size=Pt(14), align=PP_ALIGN.CENTER)
    add_textbox(s6, cx6 + Inches(0.62), cy6 + Inches(0.16), cw6 - Inches(0.78), Inches(0.28),
                title, size=Pt(12), color=GREY_LT, bold=True)
    add_textbox(s6, cx6 + Inches(0.62), cy6 + Inches(0.44), cw6 - Inches(0.78), Inches(0.22),
                label.upper(), size=Pt(8), color=accent, bold=True)
    add_textbox(s6, cx6 + Inches(0.18), cy6 + Inches(0.78), cw6 - Inches(0.30), Inches(1.4),
                desc, size=Pt(10), color=GREY_MD)
    tx = cx6 + Inches(0.18); ty_tool = cy6 + Inches(2.28)
    for t in tools:
        tw = Inches(0.12) * len(t) + Inches(0.22)
        add_rect(s6, tx, ty_tool, tw, Inches(0.30), fill=BG_DARK, line_color=BORDER)
        add_textbox(s6, tx + Inches(0.05), ty_tool + Inches(0.02), tw, Inches(0.26),
                    t, size=Pt(9), color=LIME, bold=True)
        tx += tw + Inches(0.08)
    cx6 += cw6 + Inches(0.17)

# Siemens example
ey = cy6 + ch6 + Inches(0.2); eh = Inches(1.52)
add_rect(s6, Inches(0.5), ey, Inches(12.33), eh, fill=BG_CARD, line_color=BORDER)
add_textbox(s6, Inches(0.7), ey + Inches(0.1), Inches(8), Inches(0.24),
            "WHAT ONBOARDING A BRAND LOOKS LIKE — EXAMPLE: SIEMENS", size=Pt(8), color=BLUE, bold=True)
steps = [
    ("Feed",  "SiePortal API and PDF catalogs are ingested by the KG Constructor"),
    ("Build", "A Siemens-specific knowledge graph is constructed; the product database is populated with official API declarations"),
    ("Query", "The grounding layer receives Siemens tools and answers any question without ever calling SiePortal at query time"),
]
sx = Inches(0.7)
for step, desc in steps:
    add_rect(s6, sx, ey + Inches(0.40), Inches(3.8), Inches(0.88),
             fill=RGBColor(0x16, 0x1d, 0x30), line_color=RGBColor(0x25, 0x35, 0x55))
    add_textbox(s6, sx + Inches(0.12), ey + Inches(0.44), Inches(3.55), Inches(0.26),
                step.upper(), size=Pt(9), color=BLUE, bold=True)
    add_textbox(s6, sx + Inches(0.12), ey + Inches(0.68), Inches(3.55), Inches(0.54),
                desc, size=Pt(10), color=GREY_MD)
    sx += Inches(4.08)
footer(s6, "Graph · Database · Semantic — three surfaces, one agent", "6 / 6")

# ── Save ──────────────────────────────────────────────────────────────────
prs.save("rxcore.pptx")
print("Saved: rxcore.pptx")
